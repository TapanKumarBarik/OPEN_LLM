from pptx import Presentation
from .base import BaseDocumentParser, DocumentChunk
from typing import List, Dict

class PptParser(BaseDocumentParser):
    def parse(self, file_path: str, metadata: Dict) -> List[DocumentChunk]:
        chunks = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from shapes in slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            text = " ".join(slide_text).strip()
            if not text:
                continue
                
            slide_metadata = metadata.copy()
            slide_metadata.update({
                'slide_number': slide_num,
                'total_slides': len(prs.slides)
            })
            
            slide_chunks = self._create_chunks(text, slide_metadata)
            chunks.extend(slide_chunks)
            
        return chunks